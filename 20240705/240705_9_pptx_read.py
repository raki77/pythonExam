from pptx import Presentation
prs = Presentation('presentation2.pptx')


# 각 슬라이드의 제목과 본문 텍스트를 출력
for slide_number, slide in enumerate(prs.slides): #전체 슬라이드를 순회


  print(f"Slide {slide_number+1}")

  # 슬라이드의 제목 출력 (제목이 있는 경우)
  if slide.shapes.title:
    print(f"Title: {slide.shapes.title.text}")



        
  # 슬라이드의 본문 텍스트 출력
  #각 슬라이드에 포함된 모든 개체(도형, 텍스트 상자, 표 등)를 순회
  for shape in slide.shapes: 
    # 현재 개체의 형태가 텍스트 상자인지 확인, 텍스트 상자일 경우 아래 코드를 실행
    if hasattr(shape, "text_frame") and shape.text_frame:




      for paragraph in shape.text_frame.paragraphs: # 텍스트 상자의 모든 단락을 순회
        for run in paragraph.runs: # 현재 단락의 텍스트 런에 접근
          print(run.text) #현재 텍스트 런을 출력
