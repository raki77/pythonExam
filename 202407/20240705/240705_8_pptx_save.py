# < 파워포인트 파일 생성 >

from pptx import Presentation
# Presentation 객체를 생성하여 프레젠테이션을 초기화
prs = Presentation() 
# 첫 번째 슬라이드 레이아웃 선택
slide_layout = prs.slide_layouts[0]
# 선택한 레이아웃을 사용하여 새 슬라이드 추가
slide = prs.slides.add_slide(slide_layout)



# 슬라이드에 제목과 부제목 추가
title = slide.shapes.title
subtitle = slide.placeholders[1]


# 제목과 부제목 텍스트 박스에 각각 텍스트를 추가
title.text = "Hello, World!"
subtitle.text = "파이썬 파워포인트 자동화"



# 파일 저장
prs.save('presentation2.pptx')
